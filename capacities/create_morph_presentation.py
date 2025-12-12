#!/usr/bin/env python3
"""
Create PowerPoint presentation with Morph transitions showing capacity scaling.
Uses '!! <name>' naming convention for morphing objects.
Creates shapes programmatically.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from lxml import etree

# Color palette from CLAUDE.md
FABRIC_TEAL = RGBColor(0x11, 0x78, 0x65)
FABRIC_TEAL_DARK = RGBColor(0x0D, 0x5A, 0x4C)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_blank_slide(prs):
    """Add a blank slide layout"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def set_morph_transition(slide):
    """Properly set Morph transition on a slide using XML manipulation"""
    sld = slide._element

    # Remove existing transition if present
    existing = sld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    if existing is not None:
        sld.remove(existing)

    # Create transition element with Morph
    transition = etree.SubElement(sld, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    transition.set('spd', 'med')
    transition.set('advTm', '0')

    # Add morph transition element
    morph = etree.SubElement(transition, '{http://schemas.openxmlformats.org/presentationml/2006/main}morph')
    morph.set('option', 'byObject')

def add_cu_block(slide, left, top, width, height):
    """Add a single CU block with consistent styling"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = FABRIC_TEAL
    shape.line.color.rgb = FABRIC_TEAL_DARK
    shape.line.width = Pt(1.5)
    # All CU blocks get the same name for morphing
    shape.name = "!! Capacity Visual"
    return shape

def add_text(slide, text, left, top, width, height, font_size=18):
    """Add text box to slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = RGBColor(0x32, 0x31, 0x30)
    return textbox

print("Creating presentation with Morph transitions...\n")

# SLIDE 1: Single CU with description
print("Creating Slide 1: Single CU")
slide1 = add_blank_slide(prs)

# Add single CU block in center (smaller - 1.2 inches)
cu_size = Inches(1.2)
cu_left = (prs.slide_width - cu_size) / 2
cu_top = Inches(2.5)
add_cu_block(slide1, cu_left, cu_top, cu_size, cu_size)

# Add title
add_text(slide1, "Capacity Unit (CU)", Inches(1), Inches(1), Inches(8), Inches(0.5), font_size=32)

# Add description
description = (
    "A Capacity Unit (CU) is the fundamental unit of compute in Microsoft Fabric.\n\n"
    "Each operation consumes CUs based on:\n"
    "• Complexity of the operation\n"
    "• Resources required (CPU, memory, I/O)\n"
    "• Duration of execution"
)
add_text(slide1, description, Inches(1.5), Inches(5), Inches(7), Inches(2), font_size=14)

# SLIDE 2: F64 capacity (8x8 grid - BIGGER than single CU)
print("Creating Slide 2: F64 Grid (8×8)")
slide2 = add_blank_slide(prs)
set_morph_transition(slide2)

# Add title
add_text(slide2, "F64 Capacity = 64 CUs", Inches(1), Inches(0.5), Inches(8), Inches(0.5), font_size=32)

# Create 8x8 grid - total size 4.8 inches (much bigger than 1.2")
block_size = Inches(0.5)
spacing = Inches(0.1)
grid_size = 8 * block_size + 7 * spacing  # Total grid size
start_left = (prs.slide_width - grid_size) / 2
start_top = Inches(2)

for row in range(8):
    for col in range(8):
        left = start_left + col * (block_size + spacing)
        top = start_top + row * (block_size + spacing)
        add_cu_block(slide2, left, top, block_size, block_size)

add_text(slide2, "8 × 8 = 64 Capacity Units", Inches(2.5), Inches(6.5), Inches(5), Inches(0.5), font_size=18)

# SLIDE 3: 1x64 block (one second)
print("Creating Slide 3: 1×64 vertical bar (one second)")
slide3 = add_blank_slide(prs)
set_morph_transition(slide3)

# Add title
add_text(slide3, "One Second of F64 Capacity", Inches(1), Inches(0.5), Inches(8), Inches(0.5), font_size=32)

# Create vertical stack of 64 CUs
block_width = Inches(0.4)
block_height = Inches(0.065)
stack_left = (prs.slide_width - block_width) / 2
stack_top = Inches(1.8)

for i in range(64):
    top = stack_top + i * block_height
    add_cu_block(slide3, stack_left, top, block_width, block_height)

add_text(slide3, "64 CUs stacked vertically\nrepresents 1 second of capacity",
         Inches(1), Inches(6.5), Inches(8), Inches(0.8), font_size=16)

# SLIDE 4: 30x64 grid (30 seconds)
print("Creating Slide 4: 30×64 grid (30-second window)")
slide4 = add_blank_slide(prs)
set_morph_transition(slide4)

# Add title
add_text(slide4, "30-Second Evaluation Window", Inches(1), Inches(0.5), Inches(8), Inches(0.5), font_size=32)

# Create 30x64 grid (sample every few to keep performance reasonable)
cell_width = Inches(0.24)
cell_height = Inches(0.065)
grid_start_left = Inches(0.9)
grid_start_top = Inches(1.8)

# Sample grid - show every 2nd column and every 4th row
for sec in range(0, 30, 2):
    for cu in range(0, 64, 4):
        left = grid_start_left + sec * cell_width
        top = grid_start_top + cu * cell_height
        add_cu_block(slide4, left, top, cell_width, cell_height)

add_text(slide4, "30 seconds × 64 CUs = 1,920 CU-seconds",
         Inches(1.5), Inches(6.5), Inches(7), Inches(0.5), font_size=16)

# SLIDE 5: Day view (aggregated bars)
print("Creating Slide 5: 24-hour daily overview")
slide5 = add_blank_slide(prs)
set_morph_transition(slide5)

# Add title
add_text(slide5, "24-Hour Daily Overview", Inches(1), Inches(0.5), Inches(8), Inches(0.5), font_size=32)

# Create 24 bars representing hours
bar_width = Inches(0.3)
bar_height = Inches(3.5)
bars_start_left = Inches(0.75)
bars_top = Inches(2.2)

for hour in range(24):
    left = bars_start_left + hour * (bar_width + Inches(0.03))
    add_cu_block(slide5, left, bars_top, bar_width, bar_height)

# Add time labels
add_text(slide5, "00:00", Inches(0.75), Inches(6), Inches(0.8), Inches(0.3), font_size=10)
add_text(slide5, "06:00", Inches(2.9), Inches(6), Inches(0.8), Inches(0.3), font_size=10)
add_text(slide5, "12:00", Inches(5.0), Inches(6), Inches(0.8), Inches(0.3), font_size=10)
add_text(slide5, "18:00", Inches(7.1), Inches(6), Inches(0.8), Inches(0.3), font_size=10)
add_text(slide5, "24:00", Inches(8.85), Inches(6), Inches(0.8), Inches(0.3), font_size=10)

add_text(slide5, "Aggregated view showing capacity usage across 24 hours",
         Inches(1.5), Inches(6.8), Inches(7), Inches(0.5), font_size=14)

# Save presentation
output_path = "C:\\GitHub\\conference-presentations\\capacities\\capacity-morph-demo.pptx"
prs.save(output_path)

print(f"\n✓ Presentation created successfully: {output_path}")
print("\n" + "="*60)
print("SLIDE SEQUENCE:")
print("="*60)
print("1. Single CU (1.2\") with description")
print("2. F64 capacity grid (4.8\" - 8×8) → MORPH TRANSITION")
print("3. 1×64 vertical bar (one second) → MORPH TRANSITION")
print("4. 30×64 grid (30-second window) → MORPH TRANSITION")
print("5. 24-hour daily overview → MORPH TRANSITION")
print("="*60)
print("\n✓ All CU blocks named '!! Capacity Visual' for Morph")
print("✓ F64 is 4× larger than single CU (4.8\" vs 1.2\")")
print("✓ Morph transitions properly configured via XML")
print("✓ Using Fabric Teal color scheme (#117865)")
