#!/usr/bin/env python3
"""
Fabric Capacity Metrics Presentation Generator
Creates a comprehensive PowerPoint about capacity management with focus on bursting and smoothing.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from lxml import etree
import os

# Color palette from CLAUDE.md
FABRIC_TEAL = RGBColor(0x11, 0x78, 0x65)
FABRIC_TEAL_DARK = RGBColor(0x0D, 0x5A, 0x4C)
INTERACTIVE_RED = RGBColor(0xD1, 0x34, 0x38)
BACKGROUND_BLUE = RGBColor(0x00, 0x78, 0xD4)
SUCCESS_GREEN = RGBColor(0x10, 0x7C, 0x10)
WARNING_ORANGE = RGBColor(0xA8, 0x4D, 0x0A)
POWER_BI_YELLOW = RGBColor(0xF2, 0xC8, 0x11)
TEXT_COLOR = RGBColor(0x32, 0x31, 0x30)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Presentation setup
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 widescreen
prs.slide_height = Inches(7.5)

def add_blank_slide():
    """Add a blank slide"""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)

def set_morph_transition(slide):
    """Set Morph transition on slide"""
    sld = slide._element
    existing = sld.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    if existing is not None:
        sld.remove(existing)
    transition = etree.SubElement(sld, '{http://schemas.openxmlformats.org/presentationml/2006/main}transition')
    transition.set('spd', 'med')
    morph = etree.SubElement(transition, '{http://schemas.openxmlformats.org/presentationml/2006/main}morph')
    morph.set('option', 'byObject')

def add_title(slide, text, top=Inches(0.4), font_size=40, color=TEXT_COLOR):
    """Add title text to slide"""
    textbox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.333), Inches(0.8))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = True
    return textbox

def add_subtitle(slide, text, top=Inches(1.1)):
    """Add subtitle text"""
    textbox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.333), Inches(0.5))
    tf = textbox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(20)
    p.font.color.rgb = FABRIC_TEAL_DARK
    return textbox

def add_body_text(slide, text, left=Inches(0.5), top=Inches(1.8), width=Inches(12.333), height=Inches(5), font_size=18):
    """Add body text to slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = TEXT_COLOR
    p.line_spacing = 1.5
    return textbox

def add_bullet_text(slide, bullets, left=Inches(0.7), top=Inches(1.8), width=Inches(11.5), font_size=18):
    """Add bulleted text"""
    textbox = slide.shapes.add_textbox(left, top, width, Inches(5))
    tf = textbox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_COLOR
        p.level = 0
        p.space_after = Pt(12)
    return textbox

def add_box(slide, left, top, width, height, fill_color, opacity=1.0, border_color=None, border_width=Pt(0)):
    """Add a colored box"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    return shape

def add_section_header(slide, section_num, section_title):
    """Add section header with number"""
    # Section number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6), Inches(2.5), Inches(1.333), Inches(1.333))
    circle.fill.solid()
    circle.fill.fore_color.rgb = FABRIC_TEAL
    circle.line.fill.background()

    # Number in circle
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(48)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Section title
    add_title(slide, section_title, top=Inches(4.2), font_size=36)

print("=" * 60)
print("Creating Fabric Capacity Metrics Presentation")
print("=" * 60)

# ============================================================================
# SLIDE 1: Title Slide
# ============================================================================
print("Creating Slide 1: Title")
slide = add_blank_slide()

# Title background accent
accent = add_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(2.8), FABRIC_TEAL)

# Main title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12.333), Inches(1))
tf = title_box.text_frame
p = tf.paragraphs[0]
p.text = "Microsoft Fabric Capacity Metrics"
p.font.size = Pt(48)
p.font.color.rgb = WHITE
p.font.bold = True

# Subtitle
sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(12.333), Inches(0.6))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.text = "Understanding Bursting, Smoothing & Capacity Management"
p.font.size = Pt(24)
p.font.color.rgb = WHITE

# Key topics
topics = [
    "📊  Capacity Units & SKUs",
    "⏱️  Smoothing Mechanisms (5-min & 24-hr)",
    "🚀  Bursting & Throttling",
    "🛠️  Capacity Management Strategies"
]
for i, topic in enumerate(topics):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(3.5 + i * 0.7), Inches(10), Inches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = topic
    p.font.size = Pt(22)
    p.font.color.rgb = TEXT_COLOR

# ============================================================================
# SLIDE 2: Section Header - Fundamentals
# ============================================================================
print("Creating Slide 2: Section 1 Header")
slide = add_blank_slide()
add_section_header(slide, 1, "Capacity Fundamentals")

# ============================================================================
# SLIDE 3: What is a Capacity Unit?
# ============================================================================
print("Creating Slide 3: What is a CU?")
slide = add_blank_slide()
add_title(slide, "What is a Capacity Unit (CU)?")

bullets = [
    "The fundamental measure of compute in Microsoft Fabric",
    "Every operation consumes CUs based on complexity and resources",
    "Measured in CU-seconds: 1 CU running for 1 second = 1 CU-second",
    "Different operations have different CU costs (queries, refreshes, dataflows)",
    "CU consumption is tracked and smoothed over time windows"
]
add_bullet_text(slide, bullets, top=Inches(1.6))

# Visual: Single CU block
cu_block = add_box(slide, Inches(10), Inches(2.5), Inches(1.5), Inches(1.5), FABRIC_TEAL, border_color=FABRIC_TEAL_DARK, border_width=Pt(3))
cu_block.name = "!! CU Block"
label = slide.shapes.add_textbox(Inches(10), Inches(4.2), Inches(1.5), Inches(0.4))
tf = label.text_frame
p = tf.paragraphs[0]
p.text = "1 CU"
p.font.size = Pt(14)
p.font.color.rgb = FABRIC_TEAL
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 4: SKU Overview
# ============================================================================
print("Creating Slide 4: SKU Overview")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "Fabric Capacity SKUs")
add_subtitle(slide, "From F2 to F2048 — Choose the right size for your workload")

# SKU table visualization
skus = [
    ("F2", 2), ("F4", 4), ("F8", 8), ("F16", 16), ("F32", 32),
    ("F64", 64), ("F128", 128), ("F256", 256), ("F512", 512),
    ("F1024", 1024), ("F2048", 2048)
]

# Draw SKU boxes with relative sizes
start_x = Inches(0.5)
y = Inches(2.2)
for i, (name, cus) in enumerate(skus[:6]):
    width = Inches(0.3 + cus * 0.025)
    height = Inches(0.5)
    box = add_box(slide, start_x, y, width, height, FABRIC_TEAL, border_color=FABRIC_TEAL_DARK, border_width=Pt(1))
    label = slide.shapes.add_textbox(start_x, y + Inches(0.55), width, Inches(0.3))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = f"{name}\n{cus} CU"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    start_x += width + Inches(0.15)

# Second row for larger SKUs
start_x = Inches(0.5)
y = Inches(3.8)
for name, cus in skus[6:]:
    width = Inches(0.3 + cus * 0.008)
    height = Inches(0.5)
    box = add_box(slide, start_x, y, min(width, Inches(3)), height, FABRIC_TEAL, border_color=FABRIC_TEAL_DARK, border_width=Pt(1))
    label = slide.shapes.add_textbox(start_x, y + Inches(0.55), min(width, Inches(3)), Inches(0.3))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = f"{name} ({cus} CU)"
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER
    start_x += min(width, Inches(3)) + Inches(0.15)

# Key point
note = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(12), Inches(1))
tf = note.text_frame
p = tf.paragraphs[0]
p.text = "💡 F64 is the most common starting point — 64 CUs with balanced cost/performance"
p.font.size = Pt(16)
p.font.color.rgb = FABRIC_TEAL_DARK

# ============================================================================
# SLIDE 5: The 30-Second Window
# ============================================================================
print("Creating Slide 5: 30-Second Window")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "The 30-Second Evaluation Window")
add_subtitle(slide, "How Fabric measures your capacity consumption")

bullets = [
    "Fabric evaluates capacity usage every 30 seconds",
    "F64 capacity = 64 CU × 30 seconds = 1,920 CU-seconds per window",
    "Operations are tracked and summed within each window",
    "Exceeding the window quota triggers smoothing mechanisms"
]
add_bullet_text(slide, bullets[:2], top=Inches(1.6), font_size=18)

# Visual: 30-second window grid
grid_left = Inches(1)
grid_top = Inches(3.5)
cell_w = Inches(0.35)
cell_h = Inches(0.06)

# Draw simplified grid (30 columns × sampling of rows)
for sec in range(30):
    for cu_row in range(0, 64, 8):
        box = add_box(slide,
                     grid_left + sec * cell_w,
                     grid_top + cu_row * cell_h,
                     cell_w - Inches(0.02),
                     cell_h * 7,
                     FABRIC_TEAL, border_color=FABRIC_TEAL_DARK, border_width=Pt(0.5))
        box.name = "!! Grid Cell"

# Labels
for sec in [0, 10, 20, 29]:
    label = slide.shapes.add_textbox(grid_left + sec * cell_w, grid_top + Inches(0.5), Inches(0.5), Inches(0.3))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = f"{sec}s"
    p.font.size = Pt(9)
    p.font.color.rgb = FABRIC_TEAL_DARK

calc = slide.shapes.add_textbox(Inches(1), Inches(4.3), Inches(10), Inches(0.4))
tf = calc.text_frame
p = tf.paragraphs[0]
p.text = "30 seconds × 64 CUs = 1,920 CU-seconds available per evaluation window"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_COLOR

# ============================================================================
# SLIDE 6: Section Header - Smoothing
# ============================================================================
print("Creating Slide 6: Section 2 Header")
slide = add_blank_slide()
add_section_header(slide, 2, "Smoothing Mechanisms")

# ============================================================================
# SLIDE 7: Two Types of Operations
# ============================================================================
print("Creating Slide 7: Operation Types")
slide = add_blank_slide()
add_title(slide, "Two Types of Operations")
add_subtitle(slide, "Interactive vs Background — Different smoothing rules")

# Interactive box
int_box = add_box(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.5), INTERACTIVE_RED, opacity=0.1, border_color=INTERACTIVE_RED, border_width=Pt(3))

int_title = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(5.4), Inches(0.5))
tf = int_title.text_frame
p = tf.paragraphs[0]
p.text = "🔴 Interactive Operations"
p.font.size = Pt(24)
p.font.color.rgb = INTERACTIVE_RED
p.font.bold = True

int_content = [
    "User-initiated, needs fast response",
    "Report renders, dashboard refreshes",
    "Direct queries, ad-hoc analysis",
    "",
    "Smoothing: 5 minutes",
    "Behavior: Spiky, urgent"
]
for i, line in enumerate(int_content):
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.7 + i * 0.45), Inches(5), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = line
    p.font.size = Pt(16) if i < 4 else Pt(18)
    p.font.color.rgb = TEXT_COLOR if i < 4 else INTERACTIVE_RED
    if i >= 4:
        p.font.bold = True

# Background box
bg_box = add_box(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(4.5), BACKGROUND_BLUE, opacity=0.1, border_color=BACKGROUND_BLUE, border_width=Pt(3))

bg_title = slide.shapes.add_textbox(Inches(7.2), Inches(2), Inches(5.4), Inches(0.5))
tf = bg_title.text_frame
p = tf.paragraphs[0]
p.text = "🔵 Background Operations"
p.font.size = Pt(24)
p.font.color.rgb = BACKGROUND_BLUE
p.font.bold = True

bg_content = [
    "Scheduled jobs, can wait",
    "Data refreshes, dataflows",
    "Scheduled notebooks, pipelines",
    "",
    "Smoothing: 24 hours",
    "Behavior: Stable, patient"
]
for i, line in enumerate(bg_content):
    box = slide.shapes.add_textbox(Inches(7.4), Inches(2.7 + i * 0.45), Inches(5), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = line
    p.font.size = Pt(16) if i < 4 else Pt(18)
    p.font.color.rgb = TEXT_COLOR if i < 4 else BACKGROUND_BLUE
    if i >= 4:
        p.font.bold = True

# Key insight
key = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.6))
tf = key.text_frame
p = tf.paragraphs[0]
p.text = "⚠️ This matches the Capacity Metrics App — Red = Interactive, Blue = Background"
p.font.size = Pt(16)
p.font.color.rgb = FABRIC_TEAL_DARK
p.font.bold = True

# ============================================================================
# SLIDE 8: 5-Minute Interactive Smoothing
# ============================================================================
print("Creating Slide 8: 5-Minute Smoothing")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "5-Minute Interactive Smoothing")
add_subtitle(slide, "Spreading spiky usage across a 5-minute window")

# Visual representation
chart_left = Inches(0.8)
chart_top = Inches(2)
chart_width = Inches(11)
chart_height = Inches(3)

# Background
chart_bg = add_box(slide, chart_left, chart_top, chart_width, chart_height, WHITE, border_color=FABRIC_TEAL_DARK, border_width=Pt(1))

# Capacity line
cap_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, chart_left, chart_top + Inches(0.8), chart_width, Pt(3))
cap_line.fill.solid()
cap_line.fill.fore_color.rgb = FABRIC_TEAL_DARK
cap_line.line.fill.background()

cap_label = slide.shapes.add_textbox(chart_left + chart_width + Inches(0.1), chart_top + Inches(0.65), Inches(1.2), Inches(0.3))
tf = cap_label.text_frame
p = tf.paragraphs[0]
p.text = "Capacity"
p.font.size = Pt(11)
p.font.color.rgb = FABRIC_TEAL_DARK

# Spiky raw usage bars (semi-transparent)
spikes = [0.3, 0.5, 1.5, 0.8, 2.2, 0.6, 1.8, 0.4, 1.2, 0.9]
bar_width = Inches(0.8)
for i, spike in enumerate(spikes):
    height = Inches(spike)
    bar = add_box(slide,
                 chart_left + Inches(0.3) + i * Inches(1.05),
                 chart_top + chart_height - height - Inches(0.1),
                 bar_width,
                 height,
                 INTERACTIVE_RED, border_color=INTERACTIVE_RED, border_width=Pt(1))
    bar.fill.fore_color.rgb = INTERACTIVE_RED
    # Make semi-transparent via setting alpha (simplified - just lighter color)

# Smoothed line
smooth_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     chart_left + Inches(0.3),
                                     chart_top + Inches(1.5),
                                     Inches(10.5), Pt(4))
smooth_line.fill.solid()
smooth_line.fill.fore_color.rgb = INTERACTIVE_RED
smooth_line.line.fill.background()

smooth_label = slide.shapes.add_textbox(chart_left + Inches(0.3), chart_top + Inches(1.15), Inches(2), Inches(0.3))
tf = smooth_label.text_frame
p = tf.paragraphs[0]
p.text = "5-min smoothed average"
p.font.size = Pt(11)
p.font.color.rgb = INTERACTIVE_RED
p.font.bold = True

# 5-minute window indicator
window_box = add_box(slide, chart_left + Inches(5), chart_top + chart_height + Inches(0.15), Inches(2.5), Inches(0.35), INTERACTIVE_RED, opacity=0.15)
window_label = slide.shapes.add_textbox(chart_left + Inches(5), chart_top + chart_height + Inches(0.18), Inches(2.5), Inches(0.3))
tf = window_label.text_frame
p = tf.paragraphs[0]
p.text = "← 5 minutes →"
p.font.size = Pt(12)
p.font.color.rgb = INTERACTIVE_RED
p.alignment = PP_ALIGN.CENTER

# Key insight
insight = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1))
tf = insight.text_frame
p = tf.paragraphs[0]
p.text = "A 30-second spike that uses 2× capacity gets spread over 5 minutes, reducing effective usage to a small fraction"
p.font.size = Pt(16)
p.font.color.rgb = TEXT_COLOR

# ============================================================================
# SLIDE 9: 24-Hour Background Smoothing
# ============================================================================
print("Creating Slide 9: 24-Hour Smoothing")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "24-Hour Background Smoothing")
add_subtitle(slide, "Scheduled jobs spread across the entire day")

# Visual representation
chart_left = Inches(0.8)
chart_top = Inches(2)
chart_width = Inches(11)
chart_height = Inches(3)

# Background
chart_bg = add_box(slide, chart_left, chart_top, chart_width, chart_height, WHITE, border_color=FABRIC_TEAL_DARK, border_width=Pt(1))

# Capacity line
cap_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, chart_left, chart_top + Inches(0.8), chart_width, Pt(3))
cap_line.fill.solid()
cap_line.fill.fore_color.rgb = FABRIC_TEAL_DARK
cap_line.line.fill.background()

# Scheduled job spikes
jobs = [(1, 2.0), (5, 1.8), (12, 2.5), (18, 1.5)]  # (hour position, height)
for hour, height in jobs:
    bar = add_box(slide,
                 chart_left + Inches(hour * 0.45 + 0.2),
                 chart_top + chart_height - Inches(height) - Inches(0.1),
                 Inches(0.5),
                 Inches(height),
                 BACKGROUND_BLUE)

# Smoothed line (very flat)
smooth_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                     chart_left + Inches(0.2),
                                     chart_top + Inches(2.2),
                                     Inches(10.6), Pt(4))
smooth_line.fill.solid()
smooth_line.fill.fore_color.rgb = BACKGROUND_BLUE
smooth_line.line.fill.background()

smooth_label = slide.shapes.add_textbox(chart_left + Inches(0.3), chart_top + Inches(1.9), Inches(3), Inches(0.3))
tf = smooth_label.text_frame
p = tf.paragraphs[0]
p.text = "24-hr smoothed (nearly flat!)"
p.font.size = Pt(11)
p.font.color.rgb = BACKGROUND_BLUE
p.font.bold = True

# Time labels
times = ["00:00", "06:00", "12:00", "18:00", "24:00"]
for i, time in enumerate(times):
    label = slide.shapes.add_textbox(chart_left + i * Inches(2.7), chart_top + chart_height + Inches(0.1), Inches(0.8), Inches(0.3))
    tf = label.text_frame
    p = tf.paragraphs[0]
    p.text = time
    p.font.size = Pt(10)
    p.font.color.rgb = FABRIC_TEAL_DARK

# 24-hour window indicator
window_box = add_box(slide, chart_left, chart_top + chart_height + Inches(0.4), chart_width, Inches(0.35), BACKGROUND_BLUE, opacity=0.15)
window_label = slide.shapes.add_textbox(chart_left, chart_top + chart_height + Inches(0.43), chart_width, Inches(0.3))
tf = window_label.text_frame
p = tf.paragraphs[0]
p.text = "← Full 24 hours →"
p.font.size = Pt(12)
p.font.color.rgb = BACKGROUND_BLUE
p.alignment = PP_ALIGN.CENTER

# Key insight
insight = slide.shapes.add_textbox(Inches(0.5), Inches(5.9), Inches(12), Inches(1))
tf = insight.text_frame
p = tf.paragraphs[0]
p.text = "A 1-hour job consuming 10× capacity gets spread over 24 hours = only ~40% effective utilization!"
p.font.size = Pt(16)
p.font.color.rgb = TEXT_COLOR

# ============================================================================
# SLIDE 10: Smoothing Comparison
# ============================================================================
print("Creating Slide 10: Smoothing Comparison")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "Smoothing Comparison: The Key Difference")

# Comparison table
headers = ["", "Interactive", "Background"]
rows = [
    ["Smoothing Window", "5 minutes", "24 hours"],
    ["Use Case", "User queries, reports", "Scheduled refreshes"],
    ["Spike Impact", "High (shorter window)", "Low (longer window)"],
    ["Throttle Risk", "Higher", "Lower"],
    ["Best Practice", "Optimize hot queries", "Schedule off-peak"]
]

table_left = Inches(1)
table_top = Inches(1.8)
col_widths = [Inches(3), Inches(4), Inches(4)]
row_height = Inches(0.7)

# Draw table
for col, (header, width) in enumerate(zip(headers, col_widths)):
    left = table_left + sum(w.inches for w in col_widths[:col]) * 914400 / 914400
    left = Inches(1 + col * 3.7)

    # Header row
    if col > 0:
        color = INTERACTIVE_RED if col == 1 else BACKGROUND_BLUE
        box = add_box(slide, left, table_top, Inches(3.5), row_height, color, border_color=color, border_width=Pt(2))
        label = slide.shapes.add_textbox(left, table_top + Inches(0.15), Inches(3.5), Inches(0.4))
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = header
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

# Data rows
for row_idx, row in enumerate(rows):
    for col, cell in enumerate(row):
        left = Inches(1 + col * 3.7)
        top = table_top + (row_idx + 1) * row_height

        bg_color = WHITE
        if col == 0:
            bg_color = FABRIC_TEAL
        elif col == 1:
            bg_color = RGBColor(0xFD, 0xED, 0xED)  # Light red
        else:
            bg_color = RGBColor(0xE6, 0xF2, 0xFA)  # Light blue

        box = add_box(slide, left, top, Inches(3.5), row_height, bg_color, border_color=FABRIC_TEAL_DARK, border_width=Pt(0.5))

        label = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(0.15), Inches(3.3), Inches(0.5))
        tf = label.text_frame
        p = tf.paragraphs[0]
        p.text = cell
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE if col == 0 else TEXT_COLOR
        if col == 0:
            p.font.bold = True
        p.alignment = PP_ALIGN.CENTER if col > 0 else PP_ALIGN.LEFT

# ============================================================================
# SLIDE 11: Section Header - Bursting
# ============================================================================
print("Creating Slide 11: Section 3 Header")
slide = add_blank_slide()
add_section_header(slide, 3, "Bursting Mechanism")

# ============================================================================
# SLIDE 12: What is Bursting?
# ============================================================================
print("Creating Slide 12: What is Bursting?")
slide = add_blank_slide()
add_title(slide, "What is Bursting?")
add_subtitle(slide, "Your capacity's built-in shock absorber")

bullets = [
    "Bursting allows temporary usage above your base capacity",
    "Works like a savings account: unused capacity accumulates",
    "Maximum burst bank = 10 minutes of your capacity",
    "F64 burst bank max = 64 CU × 10 min × 60 sec = 38,400 CU-seconds",
    "When you need more, you draw from the bank"
]
add_bullet_text(slide, bullets, top=Inches(1.6))

# Visual: Savings bank metaphor
bank = add_box(slide, Inches(9), Inches(2.5), Inches(3.5), Inches(2.5), POWER_BI_YELLOW, border_color=WARNING_ORANGE, border_width=Pt(3))

bank_label = slide.shapes.add_textbox(Inches(9), Inches(2.7), Inches(3.5), Inches(0.5))
tf = bank_label.text_frame
p = tf.paragraphs[0]
p.text = "🏦 Burst Bank"
p.font.size = Pt(20)
p.font.color.rgb = WARNING_ORANGE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

bank_content = slide.shapes.add_textbox(Inches(9.2), Inches(3.4), Inches(3.1), Inches(1.5))
tf = bank_content.text_frame
p = tf.paragraphs[0]
p.text = "Balance:\n38,400 CU-sec\n(10 min @ F64)"
p.font.size = Pt(16)
p.font.color.rgb = TEXT_COLOR
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 13: How Burst Accumulates
# ============================================================================
print("Creating Slide 13: Burst Accumulation")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "How Burst Capacity Accumulates")
add_subtitle(slide, "Low usage periods fill your burst bank")

# Left side: Low usage scenario
low_box = add_box(slide, Inches(0.5), Inches(1.8), Inches(5.5), Inches(4.5), SUCCESS_GREEN, opacity=0.1, border_color=SUCCESS_GREEN, border_width=Pt(2))

low_title = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(5), Inches(0.5))
tf = low_title.text_frame
p = tf.paragraphs[0]
p.text = "During Low Usage"
p.font.size = Pt(22)
p.font.color.rgb = SUCCESS_GREEN
p.font.bold = True

# Usage bars (low)
for i in range(6):
    base = add_box(slide, Inches(1 + i * 0.7), Inches(4.5), Inches(0.5), Inches(1.5), FABRIC_TEAL, opacity=0.2)
    usage = add_box(slide, Inches(1 + i * 0.7), Inches(5.5), Inches(0.5), Inches(0.5), FABRIC_TEAL)

# Arrow up
arrow_label = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(4), Inches(0.8))
tf = arrow_label.text_frame
p = tf.paragraphs[0]
p.text = "↑ Unused capacity → Burst bank"
p.font.size = Pt(16)
p.font.color.rgb = SUCCESS_GREEN

# Right side: Bank filling up
bank_box = add_box(slide, Inches(7), Inches(1.8), Inches(5.5), Inches(4.5), POWER_BI_YELLOW, opacity=0.15, border_color=WARNING_ORANGE, border_width=Pt(2))

bank_title = slide.shapes.add_textbox(Inches(7.2), Inches(2), Inches(5), Inches(0.5))
tf = bank_title.text_frame
p = tf.paragraphs[0]
p.text = "Burst Bank Fills Up"
p.font.size = Pt(22)
p.font.color.rgb = WARNING_ORANGE
p.font.bold = True

# Bank visualization (filling)
bank_container = add_box(slide, Inches(8.5), Inches(2.8), Inches(2.5), Inches(3), WHITE, border_color=WARNING_ORANGE, border_width=Pt(2))
bank_fill = add_box(slide, Inches(8.55), Inches(4.3), Inches(2.4), Inches(1.45), POWER_BI_YELLOW)

fill_label = slide.shapes.add_textbox(Inches(8.5), Inches(4.6), Inches(2.5), Inches(0.5))
tf = fill_label.text_frame
p = tf.paragraphs[0]
p.text = "~70%"
p.font.size = Pt(24)
p.font.color.rgb = WARNING_ORANGE
p.font.bold = True
p.alignment = PP_ALIGN.CENTER

max_label = slide.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(2.5), Inches(0.3))
tf = max_label.text_frame
p = tf.paragraphs[0]
p.text = "Max (10 min)"
p.font.size = Pt(10)
p.font.color.rgb = WARNING_ORANGE
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 14: How Burst Gets Consumed
# ============================================================================
print("Creating Slide 14: Burst Consumption")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "How Burst Capacity Gets Consumed")
add_subtitle(slide, "Spikes draw from your accumulated burst bank")

# Three-stage visualization
stages = [
    ("Before Spike", "Bank: 90%", SUCCESS_GREEN, 0.9),
    ("During Spike", "Using burst!", INTERACTIVE_RED, 0.6),
    ("After Spike", "Bank: 40%", WARNING_ORANGE, 0.4)
]

for i, (title, status, color, fill) in enumerate(stages):
    left = Inches(0.5 + i * 4.2)

    # Stage box
    box = add_box(slide, left, Inches(1.8), Inches(3.8), Inches(4.2), WHITE, border_color=color, border_width=Pt(2))

    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.1), Inches(1.9), Inches(3.6), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Bank container
    bank = add_box(slide, left + Inches(1), Inches(2.5), Inches(1.8), Inches(2.5), WHITE, border_color=WARNING_ORANGE, border_width=Pt(2))

    # Bank fill
    fill_height = Inches(2.4 * fill)
    bank_fill = add_box(slide, left + Inches(1.05), Inches(2.55 + Inches(2.4).inches * (1 - fill)), Inches(1.7), fill_height, POWER_BI_YELLOW)

    # Status
    status_box = slide.shapes.add_textbox(left + Inches(0.1), Inches(5.2), Inches(3.6), Inches(0.4))
    tf = status_box.text_frame
    p = tf.paragraphs[0]
    p.text = status
    p.font.size = Pt(16)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

# Arrow between stages
for i in range(2):
    arrow = slide.shapes.add_textbox(Inches(4.1 + i * 4.2), Inches(3.5), Inches(0.5), Inches(0.5))
    tf = arrow.text_frame
    p = tf.paragraphs[0]
    p.text = "→"
    p.font.size = Pt(32)
    p.font.color.rgb = FABRIC_TEAL_DARK

# Key point
key = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12), Inches(0.6))
tf = key.text_frame
p = tf.paragraphs[0]
p.text = "When burst bank empties: Interactive → Throttled delays | Background → Queued"
p.font.size = Pt(16)
p.font.color.rgb = TEXT_COLOR
p.font.bold = True

# ============================================================================
# SLIDE 15: Burst Timeline
# ============================================================================
print("Creating Slide 15: Burst Timeline")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "Burst Bank Over Time")
add_subtitle(slide, "Watching accumulation, consumption, and recovery")

# Chart area
chart_left = Inches(0.8)
chart_top = Inches(2)
chart_width = Inches(11.5)
chart_height = Inches(3.5)

chart_bg = add_box(slide, chart_left, chart_top, chart_width, chart_height, WHITE, border_color=FABRIC_TEAL_DARK, border_width=Pt(1))

# Max line
max_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, chart_left, chart_top + Inches(0.3), chart_width, Pt(2))
max_line.fill.solid()
max_line.fill.fore_color.rgb = WARNING_ORANGE
max_line.line.fill.background()

max_label = slide.shapes.add_textbox(chart_left + chart_width + Inches(0.1), chart_top + Inches(0.15), Inches(1), Inches(0.3))
tf = max_label.text_frame
p = tf.paragraphs[0]
p.text = "Max"
p.font.size = Pt(10)
p.font.color.rgb = WARNING_ORANGE

# Zero line
zero_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, chart_left, chart_top + chart_height - Inches(0.3), chart_width, Pt(2))
zero_line.fill.solid()
zero_line.fill.fore_color.rgb = FABRIC_TEAL_DARK
zero_line.line.fill.background()

# Simplified burst bank visualization with bars
# Accumulating phase
for i in range(5):
    height = Inches(0.5 + i * 0.4)
    bar = add_box(slide, chart_left + Inches(0.5 + i * 0.8), chart_top + chart_height - Inches(0.3) - height, Inches(0.6), height, POWER_BI_YELLOW)

# Spike consumption phase
for i in range(3):
    height = Inches(2.5 - i * 0.7)
    bar = add_box(slide, chart_left + Inches(4.5 + i * 0.8), chart_top + chart_height - Inches(0.3) - height, Inches(0.6), height, INTERACTIVE_RED if i > 0 else POWER_BI_YELLOW)

# Recovery phase
for i in range(5):
    height = Inches(1.1 + i * 0.3)
    bar = add_box(slide, chart_left + Inches(7 + i * 0.8), chart_top + chart_height - Inches(0.3) - height, Inches(0.6), height, POWER_BI_YELLOW)

# Phase labels
phases = [
    (Inches(1.5), "Accumulating", SUCCESS_GREEN),
    (Inches(5), "Spike!", INTERACTIVE_RED),
    (Inches(8.5), "Recovering", SUCCESS_GREEN)
]
for left, label, color in phases:
    box = slide.shapes.add_textbox(chart_left + left, chart_top + chart_height + Inches(0.1), Inches(2), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER

# Time axis
times = ["0", "5 min", "10 min", "15 min", "20 min"]
for i, time in enumerate(times):
    box = slide.shapes.add_textbox(chart_left + i * Inches(2.8), chart_top + chart_height + Inches(0.4), Inches(1), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = time
    p.font.size = Pt(10)
    p.font.color.rgb = FABRIC_TEAL_DARK

# ============================================================================
# SLIDE 16: Section Header - Throttling
# ============================================================================
print("Creating Slide 16: Section 4 Header")
slide = add_blank_slide()
add_section_header(slide, 4, "Throttling & Protection")

# ============================================================================
# SLIDE 17: Throttling States
# ============================================================================
print("Creating Slide 17: Throttling States")
slide = add_blank_slide()
add_title(slide, "Throttling States")
add_subtitle(slide, "From healthy to throttled — know the warning signs")

# Three states
states = [
    ("HEALTHY", "< 100%", "No delays", SUCCESS_GREEN),
    ("OVERAGE", "100-120%", "Using burst", WARNING_ORANGE),
    ("THROTTLED", "> 120%", "Delays active", INTERACTIVE_RED)
]

for i, (state, percent, desc, color) in enumerate(states):
    left = Inches(0.8 + i * 4.2)

    # State box
    box = add_box(slide, left, Inches(1.8), Inches(3.8), Inches(3.5), color, opacity=0.15, border_color=color, border_width=Pt(3))

    # State name
    name_box = slide.shapes.add_textbox(left, Inches(2), Inches(3.8), Inches(0.6))
    tf = name_box.text_frame
    p = tf.paragraphs[0]
    p.text = state
    p.font.size = Pt(28)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Gauge visualization
    gauge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(1.2), Inches(2.8), Inches(1.4), Inches(1.4))
    gauge.fill.solid()
    gauge.fill.fore_color.rgb = WHITE
    gauge.line.color.rgb = color
    gauge.line.width = Pt(4)

    # Percentage
    pct_box = slide.shapes.add_textbox(left, Inches(4.3), Inches(3.8), Inches(0.4))
    tf = pct_box.text_frame
    p = tf.paragraphs[0]
    p.text = percent
    p.font.size = Pt(18)
    p.font.color.rgb = color
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Description
    desc_box = slide.shapes.add_textbox(left, Inches(4.8), Inches(3.8), Inches(0.4))
    tf = desc_box.text_frame
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_COLOR
    p.alignment = PP_ALIGN.CENTER

# Behavior comparison
int_box = add_box(slide, Inches(0.8), Inches(5.6), Inches(5.8), Inches(1.4), INTERACTIVE_RED, opacity=0.1, border_color=INTERACTIVE_RED, border_width=Pt(2))
int_label = slide.shapes.add_textbox(Inches(0.9), Inches(5.7), Inches(5.6), Inches(1.2))
tf = int_label.text_frame
p = tf.paragraphs[0]
p.text = "🔴 Interactive: Requests delayed incrementally\nMax delay: ~6 minutes at 200% utilization"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_COLOR

bg_box = add_box(slide, Inches(7), Inches(5.6), Inches(5.8), Inches(1.4), BACKGROUND_BLUE, opacity=0.1, border_color=BACKGROUND_BLUE, border_width=Pt(2))
bg_label = slide.shapes.add_textbox(Inches(7.1), Inches(5.7), Inches(5.6), Inches(1.2))
tf = bg_label.text_frame
p = tf.paragraphs[0]
p.text = "🔵 Background: Jobs queued, not delayed\nMay timeout after 24 hours in queue"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_COLOR

# ============================================================================
# SLIDE 18: Throttle Delay Curve
# ============================================================================
print("Creating Slide 18: Delay Curve")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "Interactive Request Delay Curve")
add_subtitle(slide, "How delay increases with overutilization")

# Chart area
chart_left = Inches(1.5)
chart_top = Inches(2)
chart_width = Inches(10)
chart_height = Inches(4)

chart_bg = add_box(slide, chart_left, chart_top, chart_width, chart_height, WHITE, border_color=FABRIC_TEAL_DARK, border_width=Pt(1))

# Y-axis labels (Delay)
y_labels = ["6 min", "4 min", "2 min", "0"]
for i, label in enumerate(y_labels):
    box = slide.shapes.add_textbox(chart_left - Inches(0.7), chart_top + i * Inches(1.2) - Inches(0.1), Inches(0.6), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = FABRIC_TEAL_DARK
    p.alignment = PP_ALIGN.RIGHT

# X-axis labels (Utilization)
x_labels = ["0%", "50%", "100%", "150%", "200%"]
for i, label in enumerate(x_labels):
    box = slide.shapes.add_textbox(chart_left + i * Inches(2.5) - Inches(0.3), chart_top + chart_height + Inches(0.1), Inches(0.6), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = FABRIC_TEAL_DARK
    p.alignment = PP_ALIGN.CENTER

# Zone fills
# Healthy zone (0-100%)
healthy_zone = add_box(slide, chart_left, chart_top, Inches(5), chart_height, SUCCESS_GREEN, opacity=0.1)
# Overage zone (100-120%)
overage_zone = add_box(slide, chart_left + Inches(5), chart_top, Inches(1), chart_height, WARNING_ORANGE, opacity=0.1)
# Throttled zone (120%+)
throttle_zone = add_box(slide, chart_left + Inches(6), chart_top, Inches(4), chart_height, INTERACTIVE_RED, opacity=0.1)

# 100% threshold line
line_100 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, chart_left + Inches(5), chart_top, Pt(2), chart_height)
line_100.fill.solid()
line_100.fill.fore_color.rgb = FABRIC_TEAL_DARK
line_100.line.fill.background()

# 120% threshold line (throttle starts)
line_120 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, chart_left + Inches(6), chart_top, Pt(3), chart_height)
line_120.fill.solid()
line_120.fill.fore_color.rgb = INTERACTIVE_RED
line_120.line.fill.background()

# Delay curve (approximated with rectangles/bars showing increasing delay)
curve_points = [
    (Inches(0), Inches(3.9)),
    (Inches(2.5), Inches(3.9)),
    (Inches(5), Inches(3.9)),
    (Inches(6), Inches(3.7)),
    (Inches(7), Inches(2.5)),
    (Inches(8), Inches(1.5)),
    (Inches(9), Inches(0.8)),
    (Inches(10), Inches(0.3))
]

for i in range(len(curve_points) - 1):
    x1, y1 = curve_points[i]
    x2, y2 = curve_points[i + 1]
    bar = add_box(slide, chart_left + x1, chart_top + y2, x2 - x1, Pt(4), INTERACTIVE_RED)

# Zone labels
labels = [("Healthy", Inches(2), SUCCESS_GREEN), ("Overage", Inches(5.2), WARNING_ORANGE), ("Throttled", Inches(7.5), INTERACTIVE_RED)]
for label, left, color in labels:
    box = slide.shapes.add_textbox(chart_left + left, chart_top + Inches(0.2), Inches(1.5), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(12)
    p.font.color.rgb = color
    p.font.bold = True

# Axis titles
y_title = slide.shapes.add_textbox(Inches(0.5), chart_top + Inches(1.5), Inches(0.5), Inches(1))
tf = y_title.text_frame
p = tf.paragraphs[0]
p.text = "Delay"
p.font.size = Pt(12)
p.font.color.rgb = FABRIC_TEAL_DARK

x_title = slide.shapes.add_textbox(chart_left + Inches(4), chart_top + chart_height + Inches(0.5), Inches(2), Inches(0.3))
tf = x_title.text_frame
p = tf.paragraphs[0]
p.text = "Utilization %"
p.font.size = Pt(12)
p.font.color.rgb = FABRIC_TEAL_DARK

# ============================================================================
# SLIDE 19: Section Header - Management
# ============================================================================
print("Creating Slide 19: Section 5 Header")
slide = add_blank_slide()
add_section_header(slide, 5, "Capacity Management")

# ============================================================================
# SLIDE 20: Capacity Metrics App
# ============================================================================
print("Creating Slide 20: Metrics App")
slide = add_blank_slide()
add_title(slide, "Capacity Metrics App")
add_subtitle(slide, "Your visibility into capacity consumption")

bullets = [
    "Free app from Microsoft AppSource",
    "Shows CU utilization over time (last 14 days)",
    "Breaks down Interactive vs Background usage",
    "Identifies top consuming workloads and items",
    "Shows throttling events and their impact",
    "Essential for capacity planning and troubleshooting"
]
add_bullet_text(slide, bullets, left=Inches(0.7), top=Inches(1.8), width=Inches(6))

# Dashboard mockup
dash_box = add_box(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(4.5), WHITE, border_color=FABRIC_TEAL, border_width=Pt(2))

# Header bar
header = add_box(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(0.5), FABRIC_TEAL)
header_label = slide.shapes.add_textbox(Inches(7.2), Inches(1.85), Inches(5), Inches(0.4))
tf = header_label.text_frame
p = tf.paragraphs[0]
p.text = "⚡ Capacity Metrics"
p.font.size = Pt(14)
p.font.color.rgb = WHITE
p.font.bold = True

# Mini charts
chart1 = add_box(slide, Inches(7.2), Inches(2.5), Inches(2.5), Inches(1.5), RGBColor(0xF8, 0xF8, 0xF8), border_color=FABRIC_TEAL_DARK, border_width=Pt(0.5))
chart2 = add_box(slide, Inches(10), Inches(2.5), Inches(2.5), Inches(1.5), RGBColor(0xF8, 0xF8, 0xF8), border_color=FABRIC_TEAL_DARK, border_width=Pt(0.5))
chart3 = add_box(slide, Inches(7.2), Inches(4.2), Inches(5.3), Inches(1.8), RGBColor(0xF8, 0xF8, 0xF8), border_color=FABRIC_TEAL_DARK, border_width=Pt(0.5))

# Sample bars in charts
for i in range(5):
    bar = add_box(slide, Inches(7.4 + i * 0.4), Inches(3.4 - i * 0.15), Inches(0.25), Inches(0.5 + i * 0.1), FABRIC_TEAL)
    bar2 = add_box(slide, Inches(10.2 + i * 0.4), Inches(3.5), Inches(0.25), Inches(0.3), INTERACTIVE_RED)

# Interactive/Background stacked
for i in range(6):
    int_bar = add_box(slide, Inches(7.5 + i * 0.8), Inches(5.2), Inches(0.5), Inches(0.4), INTERACTIVE_RED)
    bg_bar = add_box(slide, Inches(7.5 + i * 0.8), Inches(5.6), Inches(0.5), Inches(0.3), BACKGROUND_BLUE)

# Link
link_box = slide.shapes.add_textbox(Inches(7), Inches(6.4), Inches(5.8), Inches(0.4))
tf = link_box.text_frame
p = tf.paragraphs[0]
p.text = "📥 Free in Microsoft AppSource"
p.font.size = Pt(14)
p.font.color.rgb = BACKGROUND_BLUE
p.alignment = PP_ALIGN.CENTER

# ============================================================================
# SLIDE 21: Autoscale
# ============================================================================
print("Creating Slide 21: Autoscale")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "Autoscale: Automatic Capacity Adjustment")
add_subtitle(slide, "Let Azure scale your capacity based on demand")

bullets = [
    "Automatically increases capacity when utilization is high",
    "Scales back down when demand decreases",
    "Configure maximum CU cap to control costs",
    "Helps prevent throttling during peak usage",
    "Billed only for what you use (per-second billing)"
]
add_bullet_text(slide, bullets, top=Inches(1.8), width=Inches(6))

# Autoscale visualization
chart_box = add_box(slide, Inches(6.5), Inches(2), Inches(6), Inches(3.5), WHITE, border_color=FABRIC_TEAL_DARK, border_width=Pt(1))

# Base capacity line (step function showing scale up/down)
# Base level
base1 = add_box(slide, Inches(6.7), Inches(4.5), Inches(1.5), Pt(4), FABRIC_TEAL)
# Scale up
base2 = add_box(slide, Inches(8.2), Inches(3.5), Inches(2.5), Pt(4), FABRIC_TEAL)
# Scale down
base3 = add_box(slide, Inches(10.7), Inches(4.5), Inches(1.5), Pt(4), FABRIC_TEAL)

# Vertical connectors
conn1 = add_box(slide, Inches(8.2), Inches(3.5), Pt(4), Inches(1), FABRIC_TEAL)
conn2 = add_box(slide, Inches(10.7), Inches(3.5), Pt(4), Inches(1), FABRIC_TEAL)

# Scale arrows
up_label = slide.shapes.add_textbox(Inches(7.8), Inches(3.8), Inches(0.8), Inches(0.4))
tf = up_label.text_frame
p = tf.paragraphs[0]
p.text = "↑ Scale"
p.font.size = Pt(11)
p.font.color.rgb = SUCCESS_GREEN

down_label = slide.shapes.add_textbox(Inches(10.3), Inches(4), Inches(0.8), Inches(0.4))
tf = down_label.text_frame
p = tf.paragraphs[0]
p.text = "↓ Scale"
p.font.size = Pt(11)
p.font.color.rgb = WARNING_ORANGE

# Capacity labels
labels = [("F64", Inches(7.2), Inches(4.7)), ("F128", Inches(9.2), Inches(3.2)), ("F64", Inches(11.2), Inches(4.7))]
for text, left, top in labels:
    box = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.3))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = FABRIC_TEAL

# Time axis
time_label = slide.shapes.add_textbox(Inches(9), Inches(5.6), Inches(1.5), Inches(0.3))
tf = time_label.text_frame
p = tf.paragraphs[0]
p.text = "Time →"
p.font.size = Pt(10)
p.font.color.rgb = FABRIC_TEAL_DARK

# Key point
key_box = add_box(slide, Inches(6.5), Inches(5.8), Inches(6), Inches(0.6), SUCCESS_GREEN, opacity=0.15)
key_label = slide.shapes.add_textbox(Inches(6.6), Inches(5.85), Inches(5.8), Inches(0.5))
tf = key_label.text_frame
p = tf.paragraphs[0]
p.text = "💡 Set max CU cap to prevent runaway costs"
p.font.size = Pt(14)
p.font.color.rgb = SUCCESS_GREEN

# ============================================================================
# SLIDE 22: Reservations
# ============================================================================
print("Creating Slide 22: Reservations")
slide = add_blank_slide()
set_morph_transition(slide)
add_title(slide, "Capacity Reservations")
add_subtitle(slide, "Guaranteed capacity for critical workloads")

# Left: Without reservation
left_box = add_box(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4), WHITE, border_color=INTERACTIVE_RED, border_width=Pt(2))

left_title = slide.shapes.add_textbox(Inches(0.7), Inches(2), Inches(5.4), Inches(0.4))
tf = left_title.text_frame
p = tf.paragraphs[0]
p.text = "Without Reservation"
p.font.size = Pt(20)
p.font.color.rgb = INTERACTIVE_RED
p.font.bold = True

# Shared pool
pool1 = add_box(slide, Inches(1), Inches(2.8), Inches(4.8), Inches(2), FABRIC_TEAL, opacity=0.3)
pool1_label = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(4.8), Inches(0.4))
tf = pool1_label.text_frame
p = tf.paragraphs[0]
p.text = "Shared Pool (64 CU)\nAll users compete"
p.font.size = Pt(14)
p.font.color.rgb = FABRIC_TEAL
p.alignment = PP_ALIGN.CENTER

# Users competing
users_label = slide.shapes.add_textbox(Inches(1.5), Inches(5), Inches(4), Inches(0.5))
tf = users_label.text_frame
p = tf.paragraphs[0]
p.text = "👤 👤 👤 👤  First-come, first-served"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_COLOR

# Right: With reservation
right_box = add_box(slide, Inches(7), Inches(1.8), Inches(5.8), Inches(4), WHITE, border_color=SUCCESS_GREEN, border_width=Pt(2))

right_title = slide.shapes.add_textbox(Inches(7.2), Inches(2), Inches(5.4), Inches(0.4))
tf = right_title.text_frame
p = tf.paragraphs[0]
p.text = "With Reservation"
p.font.size = Pt(20)
p.font.color.rgb = SUCCESS_GREEN
p.font.bold = True

# Reserved portion
reserved = add_box(slide, Inches(7.3), Inches(2.8), Inches(1.8), Inches(2), BACKGROUND_BLUE, opacity=0.5)
reserved_label = slide.shapes.add_textbox(Inches(7.3), Inches(3.5), Inches(1.8), Inches(0.6))
tf = reserved_label.text_frame
p = tf.paragraphs[0]
p.text = "Reserved\n16 CU"
p.font.size = Pt(12)
p.font.color.rgb = BACKGROUND_BLUE
p.alignment = PP_ALIGN.CENTER

# Shared portion
shared = add_box(slide, Inches(9.3), Inches(2.8), Inches(3), Inches(2), FABRIC_TEAL, opacity=0.3)
shared_label = slide.shapes.add_textbox(Inches(9.3), Inches(3.5), Inches(3), Inches(0.6))
tf = shared_label.text_frame
p = tf.paragraphs[0]
p.text = "Shared Pool\n48 CU"
p.font.size = Pt(12)
p.font.color.rgb = FABRIC_TEAL
p.alignment = PP_ALIGN.CENTER

# Protected user
protected_label = slide.shapes.add_textbox(Inches(7.3), Inches(5), Inches(5.4), Inches(0.5))
tf = protected_label.text_frame
p = tf.paragraphs[0]
p.text = "👤 Protected    👤 👤 👤 Shared pool"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_COLOR

# Key insight
key = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(12), Inches(0.6))
tf = key.text_frame
p = tf.paragraphs[0]
p.text = "💡 Use reservations for critical reporting workloads that must never be throttled"
p.font.size = Pt(16)
p.font.color.rgb = FABRIC_TEAL_DARK
p.font.bold = True

# ============================================================================
# SLIDE 23: Best Practices
# ============================================================================
print("Creating Slide 23: Best Practices")
slide = add_blank_slide()
add_title(slide, "Capacity Management Best Practices")

practices = [
    ("📊", "Monitor regularly", "Use Capacity Metrics App to track utilization trends"),
    ("📅", "Schedule wisely", "Run background jobs during off-peak hours"),
    ("🔍", "Optimize hot queries", "Identify and tune frequently-run interactive reports"),
    ("⚖️", "Right-size SKU", "Start with F64, scale based on actual usage"),
    ("🛡️", "Reserve for critical", "Use reservations for must-run workloads"),
    ("📈", "Enable autoscale", "Let Azure handle demand spikes automatically")
]

for i, (icon, title, desc) in enumerate(practices):
    row = i // 2
    col = i % 2
    left = Inches(0.5 + col * 6.4)
    top = Inches(1.6 + row * 1.8)

    # Icon
    icon_box = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.6))
    tf = icon_box.text_frame
    p = tf.paragraphs[0]
    p.text = icon
    p.font.size = Pt(28)

    # Title
    title_box = slide.shapes.add_textbox(left + Inches(0.7), top, Inches(5.5), Inches(0.4))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.color.rgb = FABRIC_TEAL
    p.font.bold = True

    # Description
    desc_box = slide.shapes.add_textbox(left + Inches(0.7), top + Inches(0.45), Inches(5.5), Inches(0.8))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_COLOR

# ============================================================================
# SLIDE 24: Summary
# ============================================================================
print("Creating Slide 24: Summary")
slide = add_blank_slide()
add_title(slide, "Key Takeaways")

takeaways = [
    "CUs are the fundamental measure — everything consumes capacity units",
    "30-second windows evaluate usage, but smoothing spreads the impact",
    "Interactive (5-min) vs Background (24-hr) — different smoothing, different behaviors",
    "Bursting provides a 10-minute buffer — use it wisely, don't drain it",
    "Throttling delays requests — monitor utilization to stay healthy",
    "Use Metrics App, autoscale, and reservations to manage effectively"
]

for i, takeaway in enumerate(takeaways):
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(1.6 + i * 0.9), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = FABRIC_TEAL
    circle.line.fill.background()

    num = slide.shapes.add_textbox(Inches(0.6), Inches(1.65 + i * 0.9), Inches(0.45), Inches(0.4))
    tf = num.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Takeaway text
    text_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.6 + i * 0.9), Inches(11), Inches(0.8))
    tf = text_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = takeaway
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_COLOR

# ============================================================================
# SLIDE 25: Questions / Thank You
# ============================================================================
print("Creating Slide 25: Thank You")
slide = add_blank_slide()

# Accent bar
accent = add_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(3), FABRIC_TEAL)

# Thank you text
thanks = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.333), Inches(1))
tf = thanks.text_frame
p = tf.paragraphs[0]
p.text = "Thank You!"
p.font.size = Pt(56)
p.font.color.rgb = WHITE
p.font.bold = True

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(0.6))
tf = subtitle.text_frame
p = tf.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(28)
p.font.color.rgb = WHITE

# Resources
resources_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(5), Inches(0.5))
tf = resources_title.text_frame
p = tf.paragraphs[0]
p.text = "Resources"
p.font.size = Pt(22)
p.font.color.rgb = FABRIC_TEAL
p.font.bold = True

resources = [
    "📱 Capacity Metrics App — Microsoft AppSource",
    "📖 Microsoft Fabric documentation",
    "🎓 Microsoft Learn: Fabric capacity planning"
]

for i, resource in enumerate(resources):
    box = slide.shapes.add_textbox(Inches(0.7), Inches(4.2 + i * 0.5), Inches(8), Inches(0.4))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = resource
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_COLOR

# Save presentation
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, "fabric-capacity-complete.pptx")
prs.save(output_path)

print("\n" + "=" * 60)
print("✅ PRESENTATION CREATED SUCCESSFULLY")
print("=" * 60)
print(f"📁 Output: {output_path}")
print(f"📊 Total Slides: {len(prs.slides)}")
print("\nSLIDE STRUCTURE:")
print("-" * 40)
print("1.  Title Slide")
print("2.  Section 1: Capacity Fundamentals")
print("3.  What is a Capacity Unit?")
print("4.  SKU Overview")
print("5.  The 30-Second Window")
print("6.  Section 2: Smoothing Mechanisms")
print("7.  Interactive vs Background Operations")
print("8.  5-Minute Interactive Smoothing")
print("9.  24-Hour Background Smoothing")
print("10. Smoothing Comparison")
print("11. Section 3: Bursting Mechanism")
print("12. What is Bursting?")
print("13. How Burst Accumulates")
print("14. How Burst Gets Consumed")
print("15. Burst Timeline")
print("16. Section 4: Throttling")
print("17. Throttling States")
print("18. Delay Curve")
print("19. Section 5: Capacity Management")
print("20. Capacity Metrics App")
print("21. Autoscale")
print("22. Reservations")
print("23. Best Practices")
print("24. Key Takeaways")
print("25. Thank You / Questions")
print("-" * 40)
print("\n✅ Uses Morph transitions between key slides")
print("✅ Follows CLAUDE.md color palette")
print("✅ Includes detailed bursting & smoothing visuals")
